import sys as _sys, os as _os

_SCRIPT_DIR = (
    _os.path.dirname(_os.path.abspath(__file__))
    if "__file__" in dir()
    else next(
        (p for p in _sys.path if _os.path.isfile(_os.path.join(p, "ppt_core.py"))),
        "/tmp/skills/ppt/scripts",
    )
)
if _SCRIPT_DIR not in _sys.path:
    _sys.path.insert(0, _SCRIPT_DIR)

"""
ppt_templates.py — Utilities for working with user-provided PowerPoint templates.

Usage:
    from ppt_templates import *
    prs = Presentation('/tmp/ppt_templates/corporate/template.pptx')
    clear_template_slides(prs)
    # ... build slides ...
    promote_template_visuals(prs)
    prs.save('/tmp/output.pptx')
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from PIL import Image
import copy, io


# ═══════════════════════════════════ SLIDE MANAGEMENT ═══════════════════════════════════


def clear_template_slides(prs):
    """Remove all existing slides from a template, keeping layouts intact.

    Args:
        prs: python-pptx Presentation object loaded from a template file.

    Returns:
        int: Number of slides removed.
    """
    count = 0
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        count += 1
    print(f"Cleared {count} template slides")
    return count


def populate_placeholder(
    slide,
    idx,
    text,
    font_size=None,
    font_name=None,
    font_color=None,
    bold=None,
    alignment=None,
):
    """Set text on a slide placeholder by index with optional formatting.

    Args:
        slide: python-pptx Slide object.
        idx: Placeholder index (use analyze_template() to discover indices).
        text: Text string to set.
        font_size: Optional font size in points (int).
        font_name: Optional font family name (str).
        font_color: Optional hex color string, e.g. '#1E2761'.
        bold: Optional bool for bold weight.
        alignment: Optional PP_ALIGN value (e.g. PP_ALIGN.CENTER).

    Returns:
        The placeholder shape, or None if idx not found.
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text_frame.clear()
            p = ph.text_frame.paragraphs[0]
            p.text = text
            run = p.runs[0] if p.runs else p.add_run()
            if not p.runs:
                run.text = text
            if font_size is not None:
                run.font.size = Pt(font_size)
            if font_name is not None:
                run.font.name = font_name
            if font_color is not None:
                h = font_color.lstrip("#")
                run.font.color.rgb = RGBColor(
                    int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                )
            if bold is not None:
                run.font.bold = bold
            if alignment is not None:
                p.alignment = alignment
            return ph
    print(f"  ⚠ Placeholder idx={idx} not found on this slide")
    return None


# ═══════════════════════════════════ IMAGE PLACEMENT ═══════════════════════════════════


def add_picture_constrained(slide, img_path, left_in, top_in, max_w_in, max_h_in):
    """Add a picture to a slide, constrained within a bounding box and centered.

    Unlike safe_add_picture(), this works reliably with template-based presentations
    (avoids corrupt EMU coordinate issues).

    Args:
        slide: python-pptx Slide object.
        img_path: Path to the image file.
        left_in: Left edge of bounding box in inches.
        top_in: Top edge of bounding box in inches.
        max_w_in: Maximum width in inches.
        max_h_in: Maximum height in inches.

    Returns:
        The picture shape added to the slide.
    """
    with Image.open(img_path) as img:
        iw, ih = img.size
    aspect = iw / ih

    if aspect > (max_w_in / max_h_in):
        w = max_w_in
        h = w / aspect
    else:
        h = max_h_in
        w = h * aspect

    # Center within bounding box
    x = left_in + (max_w_in - w) / 2
    y = top_in + (max_h_in - h) / 2

    return slide.shapes.add_picture(
        img_path, Inches(x), Inches(y), Inches(w), Inches(h)
    )


# ═══════════════════════════════════ BACKGROUND PROMOTION ═══════════════════════════════════


def promote_template_visuals(prs):
    """Promote inherited layout backgrounds and shapes to slide level.

    Fixes the rendering gap between the Pillow-based preview and PowerPoint.
    PowerPoint inherits blipFill backgrounds from layouts automatically;
    the preview renderer cannot. This function makes those backgrounds explicit
    by copying them as full-bleed picture shapes on each slide.

    Also copies non-placeholder shapes (logos, decorative elements) from layouts
    to their respective slides so they appear in previews.

    Run this AFTER all slides are built, BEFORE saving.

    Args:
        prs: python-pptx Presentation object.

    Returns:
        tuple: (backgrounds_promoted, shapes_promoted)
    """
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height
    fixed_bg = 0
    fixed_shapes = 0

    for i, slide in enumerate(prs.slides):
        layout = slide.slide_layout
        slide_cSld = slide._element.find(qn("p:cSld"))

        # --- Step 1: Promote blipFill background ---
        existing_bg = slide_cSld.find(qn("p:bg"))
        if existing_bg is None:
            layout_cSld = layout._element.find(qn("p:cSld"))
            layout_bg = layout_cSld.find(qn("p:bg"))
            if layout_bg is not None:
                bgPr = layout_bg.find(qn("p:bgPr"))
                if bgPr is not None:
                    blipFill = bgPr.find(qn("a:blipFill"))
                    if blipFill is not None:
                        blip = blipFill.find(qn("a:blip"))
                        rId = blip.get(qn("r:embed"))
                        if rId and rId in layout.part.rels:
                            image_part = layout.part.rels[rId].target_part
                            pic = slide.shapes.add_picture(
                                io.BytesIO(image_part.blob), 0, 0, SLIDE_W, SLIDE_H
                            )
                            # Move picture to back (behind all content)
                            spTree = slide_cSld.find(qn("p:spTree"))
                            spTree.remove(pic._element)
                            spTree.insert(2, pic._element)  # After nvGrpSpPr + grpSpPr
                            fixed_bg += 1

        # --- Step 2: Promote non-placeholder shapes (logos, decorations) ---
        for shape in layout.shapes:
            if shape.is_placeholder:
                continue
            new_shape = copy.deepcopy(shape._element)
            # Re-link any embedded images to the slide's relationship set
            for blip in new_shape.iter(qn("a:blip")):
                old_rId = blip.get(qn("r:embed"))
                if old_rId and old_rId in layout.part.rels:
                    image_part = layout.part.rels[old_rId].target_part
                    new_rId = slide.part.relate_to(image_part, RT.IMAGE)
                    blip.set(qn("r:embed"), new_rId)
            spTree = slide_cSld.find(qn("p:spTree"))
            spTree.append(new_shape)
            fixed_shapes += 1

    print(f"Promoted {fixed_bg} backgrounds, {fixed_shapes} layout shapes")
    return (fixed_bg, fixed_shapes)
