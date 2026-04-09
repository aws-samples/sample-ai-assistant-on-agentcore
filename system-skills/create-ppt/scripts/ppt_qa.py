import sys as _sys, os as _os

_SCRIPT_DIR = (
    _os.path.dirname(_os.path.abspath(__file__))
    if "__file__" in dir()
    else next(
        (p for p in _sys.path if _os.path.isfile(_os.path.join(p, "ppt_core.py"))),
        "/tmp/skills/ppt/scripts",
    )
)
if _SCRIPT_DIR not in _sys.path:
    _sys.path.insert(0, _SCRIPT_DIR)

"""
ppt_qa.py — Quality assurance: layout validation + Pillow-based slide preview.

v3: Image-background-aware contrast checking + metric-compatible font mapping.

Usage:
    from ppt_qa import *
    issues = full_audit('/tmp/output.pptx')
    paths = render_slides('/tmp/output.pptx')
    grid  = render_thumbnail_grid('/tmp/output.pptx', '/tmp/grid.png')
"""

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
import os, math


# ═══════════════════════════════════════════ FONT MAPPING ═══════════════════════════════════════════

_USER_FONTS = os.path.expanduser("~/.fonts")
_SYSTEM_FONTS = "/usr/share/fonts"

_FONT_MAP = {
    "Calibri": (
        "Carlito-Regular.ttf",
        "Carlito-Bold.ttf",
        "Carlito-Italic.ttf",
        "Carlito-BoldItalic.ttf",
    ),
    "Calibri Light": (
        "Carlito-Regular.ttf",
        "Carlito-Bold.ttf",
        "Carlito-Italic.ttf",
        "Carlito-BoldItalic.ttf",
    ),
    "Arial": (
        "LiberationSans-Regular.ttf",
        "LiberationSans-Bold.ttf",
        "LiberationSans-Italic.ttf",
        "LiberationSans-BoldItalic.ttf",
    ),
    "Arial Black": (
        "LiberationSans-Bold.ttf",
        "LiberationSans-Bold.ttf",
        "LiberationSans-BoldItalic.ttf",
        "LiberationSans-BoldItalic.ttf",
    ),
    "Helvetica": (
        "LiberationSans-Regular.ttf",
        "LiberationSans-Bold.ttf",
        "LiberationSans-Italic.ttf",
        "LiberationSans-BoldItalic.ttf",
    ),
    "Times New Roman": (
        "LiberationSerif-Regular.ttf",
        "LiberationSerif-Bold.ttf",
        "LiberationSerif-Italic.ttf",
        "LiberationSerif-BoldItalic.ttf",
    ),
    "Courier New": (
        "LiberationMono-Regular.ttf",
        "LiberationMono-Bold.ttf",
        "LiberationMono-Italic.ttf",
        "LiberationMono-BoldItalic.ttf",
    ),
    "Consolas": (
        "LiberationMono-Regular.ttf",
        "LiberationMono-Bold.ttf",
        "LiberationMono-Italic.ttf",
        "LiberationMono-BoldItalic.ttf",
    ),
    "Cambria": (
        "Caladea-Regular.ttf",
        "Caladea-Bold.ttf",
        "Caladea-Italic.ttf",
        "Caladea-BoldItalic.ttf",
    ),
    "Georgia": (
        "Lora-Regular.ttf",
        "Lora-Bold.ttf",
        "Lora-Italic.ttf",
        "Lora-BoldItalic.ttf",
    ),
    "Trebuchet MS": (
        "Carlito-Regular.ttf",
        "Carlito-Bold.ttf",
        "Carlito-Italic.ttf",
        "Carlito-BoldItalic.ttf",
    ),
    "Segoe UI": (
        "Carlito-Regular.ttf",
        "Carlito-Bold.ttf",
        "Carlito-Italic.ttf",
        "Carlito-BoldItalic.ttf",
    ),
    "Impact": (
        "LiberationSans-Bold.ttf",
        "LiberationSans-Bold.ttf",
        "LiberationSans-BoldItalic.ttf",
        "LiberationSans-BoldItalic.ttf",
    ),
    # ── Design fonts recommended by the skill (Section 0.5) ──
    "Inter": (
        "Carlito-Regular.ttf",
        "Carlito-Bold.ttf",
        "Carlito-Italic.ttf",
        "Carlito-BoldItalic.ttf",
    ),
    "DM Sans": (
        "Carlito-Regular.ttf",
        "Carlito-Bold.ttf",
        "Carlito-Italic.ttf",
        "Carlito-BoldItalic.ttf",
    ),
    "Plus Jakarta Sans": (
        "Carlito-Regular.ttf",
        "Carlito-Bold.ttf",
        "Carlito-Italic.ttf",
        "Carlito-BoldItalic.ttf",
    ),
    "IBM Plex Sans": (
        "LiberationSans-Regular.ttf",
        "LiberationSans-Bold.ttf",
        "LiberationSans-Italic.ttf",
        "LiberationSans-BoldItalic.ttf",
    ),
}

_FALLBACK_ENTRY = (
    "DroidSans.ttf",
    "DroidSans-Bold.ttf",
    "DroidSans.ttf",
    "DroidSans-Bold.ttf",
)
_FONT_CACHE = {}


def _resolve_font_path(font_name, bold=False, italic=False):
    """Resolve a PPTX font name + style to a local TTF file path."""
    entry = _FONT_MAP.get(font_name)
    if entry is None:
        entry = _FONT_MAP.get("Arial", _FALLBACK_ENTRY)
    idx = (1 if bold else 0) + (2 if italic else 0)
    filename = entry[idx]
    for base in [_USER_FONTS, _SYSTEM_FONTS + "/google-droid-sans-fonts"]:
        path = os.path.join(base, filename)
        if os.path.exists(path):
            return path
    for base in [_USER_FONTS, _SYSTEM_FONTS]:
        for root, dirs, files in os.walk(base):
            if filename in files:
                return os.path.join(root, filename)
    return _SYSTEM_FONTS + "/google-droid-sans-fonts/DroidSans.ttf"


def _get_pil_font(font_name, bold, italic, size_px):
    """Get a cached PIL ImageFont for the given PPTX font spec."""
    key = (font_name, bold, italic, size_px)
    if key not in _FONT_CACHE:
        path = _resolve_font_path(font_name, bold, italic)
        try:
            _FONT_CACHE[key] = ImageFont.truetype(path, size_px)
        except:
            _FONT_CACHE[key] = ImageFont.load_default()
    return _FONT_CACHE[key]


# ═══════════════════════════════════════════ TEXT WRAPPING ═══════════════════════════════════════════


def _wrap_text(text, font, max_width_px):
    """Word-wrap text to fit within max_width_px. Returns list of lines."""
    if max_width_px <= 0:
        return [text]
    result_lines = []
    for paragraph in text.split("\n"):
        if not paragraph:
            result_lines.append("")
            continue
        words = paragraph.split(" ")
        current_line = ""
        for word in words:
            test = (current_line + " " + word).strip()
            if font.getlength(test) <= max_width_px or not current_line:
                current_line = test
            else:
                result_lines.append(current_line)
                current_line = word
                if font.getlength(word) > max_width_px:
                    result_lines.pop()
                    char_line = ""
                    for ch in word:
                        test_ch = char_line + ch
                        if font.getlength(test_ch) > max_width_px and char_line:
                            result_lines.append(char_line)
                            char_line = ch
                        else:
                            char_line = test_ch
                    current_line = char_line
        if current_line:
            result_lines.append(current_line)
    return result_lines if result_lines else [""]


# ═══════════════════════════════════════════ VALIDATION ═══════════════════════════════════════════


def validate_layout(filepath, safe_margin_inches=0.3):
    prs = Presentation(filepath)
    slide_w, slide_h = prs.slide_width, prs.slide_height
    tol = Inches(safe_margin_inches)
    issues = []
    for si, slide in enumerate(prs.slides, 1):
        shapes = []
        for sh in slide.shapes:
            l, t, w, h = (
                (sh.left or 0),
                (sh.top or 0),
                (sh.width or 0),
                (sh.height or 0),
            )
            r, b = l + w, t + h
            name = sh.name or "Unnamed"
            shapes.append(
                {
                    "name": name,
                    "left": l,
                    "top": t,
                    "right": r,
                    "bottom": b,
                    "width": w,
                    "height": h,
                    "is_picture": sh.shape_type == MSO_SHAPE_TYPE.PICTURE,
                    "is_textbox": sh.has_text_frame,
                }
            )
            if r > slide_w + tol:
                issues.append(
                    {
                        "slide": si,
                        "shape": name,
                        "severity": "ERROR",
                        "issue": "OVERFLOW_RIGHT",
                        "detail": f"right={r / 914400:.1f}in",
                    }
                )
            if b > slide_h + tol:
                issues.append(
                    {
                        "slide": si,
                        "shape": name,
                        "severity": "ERROR",
                        "issue": "OVERFLOW_BOTTOM",
                        "detail": f"bottom={b / 914400:.1f}in",
                    }
                )
            if l < -tol:
                issues.append(
                    {
                        "slide": si,
                        "shape": name,
                        "severity": "ERROR",
                        "issue": "OVERFLOW_LEFT",
                    }
                )
            if t < -tol:
                issues.append(
                    {
                        "slide": si,
                        "shape": name,
                        "severity": "ERROR",
                        "issue": "OVERFLOW_TOP",
                    }
                )
        pics = [s for s in shapes if s["is_picture"] and s["width"] > Inches(2)]
        txts = [s for s in shapes if s["is_textbox"] and s["width"] > Inches(1)]
        for p in pics:
            for tb in txts:
                if (
                    p["left"] < tb["right"]
                    and p["right"] > tb["left"]
                    and p["top"] < tb["bottom"]
                    and p["bottom"] > tb["top"]
                ):
                    ox = min(p["right"], tb["right"]) - max(p["left"], tb["left"])
                    oy = min(p["bottom"], tb["bottom"]) - max(p["top"], tb["top"])
                    sq = (ox / 914400) * (oy / 914400)
                    if sq > 1.0:
                        issues.append(
                            {
                                "slide": si,
                                "severity": "WARNING",
                                "shape": f"'{p['name']}' & '{tb['name']}'",
                                "issue": "OVERLAP",
                                "detail": f"~{sq:.1f} sq in",
                            }
                        )
    return issues


def check_font_compliance(filepath, allowed_fonts=None):
    if allowed_fonts is None:
        allowed_fonts = {
            "Calibri",
            "Arial",
            "Helvetica",
            "Segoe UI",
            "Calibri Light",
            "Arial Black",
            "Georgia",
            "Cambria",
            "Trebuchet MS",
            "Impact",
            "Palatino",
            "Garamond",
            "Consolas",
            # Design fonts recommended by the skill (Section 0.5)
            "Inter",
            "DM Sans",
            "Plus Jakarta Sans",
            "IBM Plex Sans",
        }
    prs = Presentation(filepath)
    violations = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name and run.font.name not in allowed_fonts:
                            violations.append(
                                {
                                    "slide": si,
                                    "shape": shape.name,
                                    "severity": "WARNING",
                                    "issue": "FONT_NOT_ALLOWED",
                                    "detail": f"Font: {run.font.name}",
                                }
                            )
    return violations


def check_min_font_size(filepath, min_pt=14):
    prs = Presentation(filepath)
    violations = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (
                            run.font.size
                            and run.font.size.pt < min_pt
                            and run.text.strip()
                        ):
                            violations.append(
                                {
                                    "slide": si,
                                    "shape": shape.name,
                                    "severity": "WARNING",
                                    "issue": "FONT_TOO_SMALL",
                                    "detail": f'{run.font.size.pt}pt < {min_pt}pt, text: "{run.text[:30]}"',
                                }
                            )
    return violations


def _relative_luminance(hex_color):
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255

    def lin(c):
        return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4

    return 0.2126 * lin(r) + 0.7152 * lin(g) + 0.0722 * lin(b)


def _contrast_ratio(hex1, hex2):
    l1, l2 = _relative_luminance(hex1), _relative_luminance(hex2)
    return (max(l1, l2) + 0.05) / (min(l1, l2) + 0.05)


def _has_full_bleed_image(slide, slide_w, slide_h):
    """Return True if the slide has a picture covering ≥85% of its area."""
    slide_area = slide_w * slide_h
    if slide_area <= 0:
        return False
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            w = shape.width or 0
            h = shape.height or 0
            if (w * h) >= slide_area * 0.85:
                return True
    return False


def _find_effective_bg(shape, all_shapes, shape_fills, slide_bg):
    own_fill = shape_fills.get(shape.name)
    if own_fill:
        return own_fill
    sl = shape.left or 0
    st = shape.top or 0
    sr = sl + (shape.width or 0)
    sb = st + (shape.height or 0)
    tol = int(0.05 * 914400)
    best_fill = None
    best_area = float("inf")
    for other_shape, other_name, other_fill in all_shapes:
        if other_shape is shape:
            continue
        if other_fill is None:
            continue
        ol = other_shape.left or 0
        ot = other_shape.top or 0
        ow = other_shape.width or 0
        oh = other_shape.height or 0
        oright = ol + ow
        ob = ot + oh
        if ol <= sl + tol and ot <= st + tol and oright >= sr - tol and ob >= sb - tol:
            area = ow * oh
            if area < best_area:
                best_area = area
                best_fill = other_fill
    return best_fill or slide_bg


def check_contrast_ratios(filepath, min_ratio_normal=4.5, min_ratio_large=3.0):
    prs = Presentation(filepath)
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height
    violations = []
    for si, slide in enumerate(prs.slides, 1):
        # Detect effective slide background
        has_image_bg = _has_full_bleed_image(slide, slide_w_emu, slide_h_emu)
        slide_bg = None  # None = unknown (image-based)
        if not has_image_bg:
            slide_bg = "#FFFFFF"
            try:
                bg_fill = slide.background.fill
                if bg_fill.type is not None:
                    slide_bg = f"#{bg_fill.fore_color.rgb}"
            except:
                pass

        shape_fills = {}
        shape_list = []
        for shape in slide.shapes:
            fill_hex = None
            try:
                if shape.fill.type is not None:
                    fill_hex = f"#{shape.fill.fore_color.rgb}"
                    shape_fills[shape.name] = fill_hex
            except:
                pass
            shape_list.append((shape, shape.name, fill_hex))

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            bg_color = _find_effective_bg(shape, shape_list, shape_fills, slide_bg)
            # Skip contrast check if effective bg is unknown (image-based)
            if bg_color is None:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    try:
                        if run.font.color.type is not None and run.font.color.rgb:
                            text_color = f"#{run.font.color.rgb}"
                        else:
                            continue
                    except:
                        continue
                    size_pt = run.font.size.pt if run.font.size else 14
                    is_bold = run.font.bold or False
                    is_large = size_pt >= 18 or (size_pt >= 14 and is_bold)
                    try:
                        ratio = _contrast_ratio(text_color, bg_color)
                        min_req = min_ratio_large if is_large else min_ratio_normal
                        if ratio < min_req:
                            violations.append(
                                {
                                    "slide": si,
                                    "shape": shape.name,
                                    "severity": "WARNING",
                                    "issue": "LOW_CONTRAST",
                                    "detail": f'{ratio:.1f}:1 (need {min_req}:1) | text={text_color} bg={bg_color} | "{run.text[:25]}"',
                                }
                            )
                    except:
                        pass
    return violations


def check_content_zone(filepath, header_bottom=1.70, footer_top=7.0):
    prs = Presentation(filepath)
    violations = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            l, t = (shape.left or 0) / 914400, (shape.top or 0) / 914400
            w, h = (shape.width or 0) / 914400, (shape.height or 0) / 914400
            b = t + h
            if w > 12 and h > 6:
                continue
            if w < 0.3 and h < 0.3:
                continue
            if t < header_bottom and b > header_bottom and h > 0.5:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    violations.append(
                        {
                            "slide": si,
                            "shape": shape.name,
                            "severity": "INFO",
                            "issue": "HEADER_INTRUSION",
                            "detail": f"top={t:.2f}in crosses header_bottom={header_bottom}in",
                        }
                    )
    return violations


def full_audit(filepath, print_report=True):
    all_issues = []
    all_issues.extend(validate_layout(filepath))
    all_issues.extend(check_font_compliance(filepath))
    all_issues.extend(check_min_font_size(filepath))
    all_issues.extend(check_contrast_ratios(filepath))
    all_issues.extend(check_content_zone(filepath))
    if print_report:
        if not all_issues:
            print("AUDIT PASSED — no issues found")
        else:
            errors = [i for i in all_issues if i["severity"] == "ERROR"]
            warnings = [i for i in all_issues if i["severity"] == "WARNING"]
            infos = [i for i in all_issues if i["severity"] == "INFO"]
            print(
                f"AUDIT: {len(errors)} errors, {len(warnings)} warnings, {len(infos)} info"
            )
            for i in all_issues:
                print(
                    f"  [{i['severity']}] Slide {i['slide']} | {i['issue']} | {i.get('shape', '')} | {i.get('detail', '')}"
                )
    return all_issues


# ═══════════════════════════════════════════ RENDERING ═══════════════════════════════════════════


def _pptx_color_to_rgb(color_obj):
    try:
        if color_obj and color_obj.rgb:
            h = str(color_obj.rgb)
            return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except:
        pass
    return None


def _get_fill_color(shape):
    try:
        if shape.fill.type is not None:
            return _pptx_color_to_rgb(shape.fill.fore_color)
    except:
        pass
    return None


def _get_slide_bg_color(slide):
    try:
        bg = slide.background.fill
        if bg.type is not None:
            return _pptx_color_to_rgb(bg.fore_color)
    except:
        pass
    return (255, 255, 255)


def _get_text_box_margins(shape, dpi):
    """Extract text box internal margins in pixels."""
    tf = shape.text_frame
    DEFAULT_LR = int(0.1 * 914400)
    DEFAULT_TB = int(0.05 * 914400)
    ml = (tf.margin_left if tf.margin_left is not None else DEFAULT_LR) / 914400 * dpi
    mr = (tf.margin_right if tf.margin_right is not None else DEFAULT_LR) / 914400 * dpi
    mt = (tf.margin_top if tf.margin_top is not None else DEFAULT_TB) / 914400 * dpi
    mb = (
        (tf.margin_bottom if tf.margin_bottom is not None else DEFAULT_TB)
        / 914400
        * dpi
    )
    return int(ml), int(mt), int(mr), int(mb)


def _get_alignment(para):
    """Get paragraph alignment, defaulting to LEFT."""
    try:
        if para.alignment == PP_ALIGN.CENTER:
            return "center"
        elif para.alignment == PP_ALIGN.RIGHT:
            return "right"
    except:
        pass
    return "left"


def render_slides(filepath, output_dir="/tmp/slide_previews", dpi=150):
    """Render all slides to PNG with accurate font mapping and text wrapping."""
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(filepath)
    sw_emu, sh_emu = prs.slide_width, prs.slide_height
    sw_px = int(sw_emu / 914400 * dpi)
    sh_px = int(sh_emu / 914400 * dpi)
    image_paths = []

    for slide_idx, slide in enumerate(prs.slides):
        bg_color = _get_slide_bg_color(slide)
        img = Image.new("RGB", (sw_px, sh_px), bg_color)
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            l = int((shape.left or 0) / 914400 * dpi)
            t = int((shape.top or 0) / 914400 * dpi)
            w = int((shape.width or 0) / 914400 * dpi)
            h = int((shape.height or 0) / 914400 * dpi)
            if w <= 0 or h <= 0:
                continue

            fill_color = _get_fill_color(shape)
            if fill_color:
                is_rounded = False
                try:
                    if hasattr(shape, "_element") and "roundRect" in shape._element.xml:
                        is_rounded = True
                except:
                    pass
                if is_rounded:
                    draw.rounded_rectangle(
                        [l, t, l + w, t + h], radius=min(w, h) // 8, fill=fill_color
                    )
                else:
                    draw.rectangle([l, t, l + w, t + h], fill=fill_color)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    import io

                    pic = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                    pic = pic.resize((w, h), Image.LANCZOS)
                    img.paste(pic, (l, t))
                except:
                    draw.rectangle(
                        [l, t, l + w, t + h],
                        fill=(220, 220, 220),
                        outline=(180, 180, 180),
                    )
                    draw.line([(l, t), (l + w, t + h)], fill=(180, 180, 180))
                    draw.line([(l + w, t), (l, t + h)], fill=(180, 180, 180))

            if shape.has_text_frame:
                ml, mt_m, mr, mb = _get_text_box_margins(shape, dpi)
                text_area_w = w - ml - mr
                cursor_y = t + mt_m

                for para in shape.text_frame.paragraphs:
                    alignment = _get_alignment(para)
                    para_runs = [(run.text, run.font) for run in para.runs if run.text]
                    if not para_runs:
                        cursor_y += int(14 * dpi / 72 * 1.2)
                        continue

                    for run_text, run_font in para_runs:
                        if not run_text.strip() and not run_text:
                            continue
                        font_name = run_font.name or "Calibri"
                        font_size_pt = run_font.size.pt if run_font.size else 14
                        is_bold = run_font.bold or False
                        is_italic = run_font.italic or False
                        text_color = (51, 51, 51)
                        try:
                            c = _pptx_color_to_rgb(run_font.color)
                            if c:
                                text_color = c
                        except:
                            pass

                        font_size_px = int(font_size_pt * dpi / 72)
                        pil_font = _get_pil_font(
                            font_name, is_bold, is_italic, font_size_px
                        )
                        line_height = int(font_size_px * 1.2)
                        wrapped_lines = _wrap_text(run_text, pil_font, text_area_w)

                        for line in wrapped_lines:
                            if not line and not run_text:
                                cursor_y += line_height
                                continue
                            line_w = pil_font.getlength(line) if line.strip() else 0
                            if alignment == "center":
                                text_x = l + ml + int((text_area_w - line_w) / 2)
                            elif alignment == "right":
                                text_x = l + ml + int(text_area_w - line_w)
                            else:
                                text_x = l + ml
                            draw.text(
                                (text_x, cursor_y), line, fill=text_color, font=pil_font
                            )
                            cursor_y += line_height
                    cursor_y += int(0.02 * dpi)

        out_path = os.path.join(output_dir, f"slide_{slide_idx + 1:02d}.png")
        img.save(out_path, "PNG")
        image_paths.append(out_path)
    return image_paths


def render_thumbnail_grid(
    filepath,
    output_path="/tmp/slide_grid.png",
    cols=2,
    thumb_w=600,
    padding=20,
    dpi=150,
):
    """Render all slides and arrange as a thumbnail grid."""
    slide_paths = render_slides(filepath, dpi=dpi)
    if not slide_paths:
        return None
    thumbs = []
    for p in slide_paths:
        im = Image.open(p)
        aspect = im.height / im.width
        thumbs.append(im.resize((thumb_w, int(thumb_w * aspect)), Image.LANCZOS))
    rows = math.ceil(len(thumbs) / cols)
    thumb_h = thumbs[0].height
    label_h = 25
    grid_w = cols * thumb_w + (cols + 1) * padding
    grid_h = rows * (thumb_h + label_h) + (rows + 1) * padding
    grid = Image.new("RGB", (grid_w, grid_h), (240, 240, 240))
    draw = ImageDraw.Draw(grid)
    try:
        font = ImageFont.truetype(_resolve_font_path("Arial", False, False), 16)
    except:
        font = ImageFont.load_default()
    for idx, thumb in enumerate(thumbs):
        row, col = divmod(idx, cols)
        x = padding + col * (thumb_w + padding)
        y = padding + row * (thumb_h + label_h + padding)
        draw.rectangle(
            [x + 3, y + 3, x + thumb_w + 3, y + thumb_h + 3], fill=(200, 200, 200)
        )
        grid.paste(thumb, (x, y))
        draw.rectangle(
            [x, y, x + thumb_w, y + thumb_h], outline=(180, 180, 180), width=1
        )
        draw.text(
            (x, y + thumb_h + 4), f"Slide {idx + 1}", fill=(80, 80, 80), font=font
        )
    grid.save(output_path, "PNG", quality=95)
    return output_path
