import sys as _sys, os as _os

_SCRIPT_DIR = (
    _os.path.dirname(_os.path.abspath(__file__))
    if "__file__" in dir()
    else next(
        (p for p in _sys.path if _os.path.isfile(_os.path.join(p, "ppt_core.py"))),
        "/tmp/skills/ppt/scripts",
    )
)
if _SCRIPT_DIR not in _sys.path:
    _sys.path.insert(0, _SCRIPT_DIR)

"""
ppt_core.py — Foundation layer for PPT generation.

Provides standard constants, color palette setup, and helper functions
that all other ppt_* modules depend on.

IMPORTANT — COORDINATE CONVENTION:
    All position and size parameters (x, y, w, h, left, top, etc.) across
    every function in this module accept **raw float values in inches**.
    They are wrapped with Inches() internally — do NOT pass Inches() objects.

    ✓  add_text(slide, 0.8, 1.7, 11.0, 0.5, "Hello")
    ✗  add_text(slide, Inches(0.8), Inches(1.7), ...)   ← WRONG, double-wraps

Usage:
    from ppt_core import *
    prs = init_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, 1.0, TITLE_Y, 5.0, 0.5, "Hello", size=TITLE_SIZE, bold=True, color=PALETTE['primary'])
    card(slide, 1.0, CONTENT_Y, 5.0, 3.0)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import copy

# ── Slide dimensions (widescreen 16:9) ──────────────────
SLIDE_W = 13.333
SLIDE_H = 7.5

# ── Standard header anchors ─────────────────────────────
BADGE_Y = 0.55
TITLE_Y = 0.95
TITLE_SIZE = 28
CONTENT_Y = 1.70

# ── Content zone boundaries ─────────────────────────────
CONTENT_BOTTOM = 7.0
CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_Y  # 5.30"
MARGIN_LEFT = 0.8
MARGIN_RIGHT = 0.8
BODY_WIDTH = SLIDE_W - MARGIN_LEFT - MARGIN_RIGHT  # ~11.73"

# ── Two-column layout anchors ───────────────────────────
LEFT_COL_X = 0.8
LEFT_COL_MAX_W = 5.5  # ends at x=6.3
GUTTER_LEFT = 6.3
GUTTER_RIGHT = 6.8
RIGHT_COL_X = 6.8
RIGHT_COL_MAX_W = SLIDE_W - RIGHT_COL_X - MARGIN_RIGHT  # ~5.73"

# ── Three-column layout anchors ─────────────────────────
COL3_POSITIONS = [
    (0.8, 3.7),  # col 1: x=0.8,  w=3.7
    (4.9, 3.7),  # col 2: x=4.9,  w=3.7
    (9.0, 3.53),  # col 3: x=9.0,  w=3.53
]

# ── Default palette (override per project) ──────────────
PALETTE = {
    "primary": "#1B2A4A",
    "secondary": "#2D5F8A",
    "accent": "#22C55E",
    "text_dark": "#333333",
    "text_body": "#4A4A4A",
    "text_muted": "#888888",
    "bg_white": "#FFFFFF",
    "bg_light": "#F5F5F5",
    "card_fill": "#F0F5FB",
    "card_fill_alt": "#E8F0FE",
    "white": "#FFFFFF",
    "series": ["#1B2A4A", "#2D5F8A", "#22C55E", "#F59E0B", "#EF4444", "#8B5CF6"],
}


def hex_to_rgb(hex_str):
    """Convert '#RRGGBB' or 'RRGGBB' to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _relative_luminance(hex_color):
    """Relative luminance of a hex color (0.0 = black, 1.0 = white)."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255

    def lin(c):
        return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4

    return 0.2126 * lin(r) + 0.7152 * lin(g) + 0.0722 * lin(b)


def _contrast_ratio(hex1, hex2):
    """WCAG contrast ratio between two hex colors."""
    l1, l2 = _relative_luminance(hex1), _relative_luminance(hex2)
    return (max(l1, l2) + 0.05) / (min(l1, l2) + 0.05)


def set_palette(overrides: dict):
    """Update PALETTE with project-specific colors."""
    PALETTE.update(overrides)


def validate_palette_contrast(print_report=True):
    """Check that palette colors have adequate contrast for their typical uses.

    Tests accent, text_muted, text_body, and primary against card fills and
    white backgrounds using WCAG AA thresholds (4.5:1 normal, 3.0:1 large).

    Returns list of issue dicts.
    """
    issues = []
    checks = [
        ("accent", "card_fill", "KPI deltas / callout numbers on cards", True, 3.0),
        ("accent", "card_fill_alt", "accent text on alt cards", True, 3.0),
        ("accent", "bg_white", "accent text on white background", True, 3.0),
        ("text_muted", "bg_white", "muted labels on white", False, 4.5),
        ("text_muted", "card_fill", "muted labels on cards", False, 4.5),
        ("text_body", "card_fill", "body text on cards", False, 4.5),
        ("text_body", "bg_white", "body text on white", False, 4.5),
        ("primary", "card_fill", "primary headings on cards", True, 3.0),
        ("primary", "bg_white", "primary headings on white", True, 3.0),
    ]
    for color_key, bg_key, usage, is_large, min_ratio in checks:
        fg = PALETTE.get(color_key)
        bg = PALETTE.get(bg_key)
        if fg is None or bg is None:
            continue
        ratio = _contrast_ratio(fg, bg)
        if ratio < min_ratio:
            issues.append(
                {
                    "text_color": f"{color_key} ({fg})",
                    "bg_color": f"{bg_key} ({bg})",
                    "ratio": round(ratio, 2),
                    "required": min_ratio,
                    "usage": usage,
                }
            )
    if print_report:
        if not issues:
            print("PALETTE CONTRAST: all checks passed ✓")
        else:
            print(f"PALETTE CONTRAST: {len(issues)} problem(s) found")
            for i in issues:
                print(
                    f"  ✗ {i['text_color']} on {i['bg_color']}: "
                    f"{i['ratio']}:1 (need {i['required']}:1) — {i['usage']}"
                )
    return issues


def init_presentation(template_path=None):
    """Create a widescreen 16:9 presentation."""
    prs = Presentation(template_path) if template_path else Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def blank_slide(prs):
    """Add a blank slide (layout index 6) and return it."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _resolve_align(align):
    """Normalize alignment to PP_ALIGN enum. Accepts strings or PP_ALIGN values."""
    if isinstance(align, str):
        return {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }.get(align.lower(), PP_ALIGN.LEFT)
    # Accept PP_ALIGN enum values directly
    if hasattr(align, "value"):
        return align
    return PP_ALIGN.LEFT


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=16,
    bold=False,
    italic=False,
    color=None,
    font_name="Calibri",
    align="left",
    valign="top",
    word_wrap=True,
    line_spacing=1.05,
    # ── Aliases (accepted but not advertised as primary) ──
    font_size=None,
    alignment=None,
):
    """Add a text box to a slide. Returns the textbox shape.

    All position/size params are raw floats in inches (wrapped internally).

    Args:
        slide:        Slide object
        x, y:         Top-left position in inches (float, NOT Inches())
        w, h:         Width and height in inches (float, NOT Inches())
        text:         String content
        size:         Font size in points (number, not Pt()). Alias: font_size
        bold:         True for bold weight
        italic:       True for italic style
        color:        Hex string '#RRGGBB' or RGBColor
        font_name:    Font family name (str)
        align:        'left', 'center', 'right' (or PP_ALIGN enum). Alias: alignment
        valign:       'top', 'middle', 'bottom'
        word_wrap:    Enable word wrapping
        line_spacing: Line spacing multiplier (e.g. 1.05)
    """
    # ── Resolve aliases ──
    if font_size is not None:
        size = font_size
    if alignment is not None:
        align = alignment
    # Handle Pt() being passed for size — extract the raw number
    if hasattr(size, "pt"):
        size = size.pt

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    tf.paragraphs[0].alignment = _resolve_align(align)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name

    if color is None:
        color = PALETTE["text_dark"]
    if isinstance(color, str):
        run.font.color.rgb = hex_to_rgb(color)
    elif isinstance(color, RGBColor):
        run.font.color.rgb = color
    else:
        run.font.color.rgb = hex_to_rgb(PALETTE["text_dark"])

    from pptx.oxml.ns import qn

    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
    spcPct = lnSpc.makeelement(qn("a:spcPct"), {"val": str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    return txBox


def add_multiline_text(
    slide,
    x,
    y,
    w,
    h,
    lines,
    size=16,
    bold=False,
    color=None,
    font_name="Calibri",
    line_spacing=1.15,
    space_after_pt=6,
    # ── Aliases ──
    font_size=None,
):
    """Add a text box with multiple paragraphs.

    All position/size params are raw floats in inches (wrapped internally).
    `lines` is a list of strings or dicts with keys: text, size, bold, color.

    Args:
        slide:          Slide object
        x, y:           Top-left position in inches (float, NOT Inches())
        w, h:           Width and height in inches (float, NOT Inches())
        lines:          List of str or dict(text, size, bold, color)
        size:           Default font size in points (number). Alias: font_size
        bold:           Default bold setting
        color:          Default hex color string
        font_name:      Font family name
        line_spacing:   Line spacing multiplier
        space_after_pt: Space after each paragraph in points
    """
    if font_size is not None:
        size = font_size
    if hasattr(size, "pt"):
        size = size.pt

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    default_color = color or PALETTE["text_body"]

    for i, line in enumerate(lines):
        if isinstance(line, str):
            line = {"text": line}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line.get("text", "")
        _s = line.get("size", size)
        if hasattr(_s, "pt"):
            _s = _s.pt
        run.font.size = Pt(_s)
        run.font.bold = line.get("bold", bold)
        run.font.name = font_name
        c = line.get("color", default_color)
        if isinstance(c, str):
            run.font.color.rgb = hex_to_rgb(c)
        else:
            run.font.color.rgb = c
        p.space_after = Pt(line.get("space_after", space_after_pt))
    return txBox


def card(
    slide,
    x,
    y,
    w,
    h,
    fill_color=None,
    corner_radius=0.15,
    border=False,
    border_color=None,
    border_width=0.75,
):
    """Add a modern card (rounded rectangle). Returns the shape.

    All position/size params are raw floats in inches (wrapped internally).

    Args:
        slide:         Slide object
        x, y:          Top-left position in inches (float, NOT Inches())
        w, h:          Width and height in inches (float, NOT Inches())
        fill_color:    Hex string or RGBColor (default: PALETTE['card_fill'])
        corner_radius: Corner radius (not currently adjustable via python-pptx)
        border:        True to add border line
        border_color:  Hex string for border (default: '#D0D0D0')
        border_width:  Border width in points
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    fc = fill_color or PALETTE["card_fill"]
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fc) if isinstance(fc, str) else fc
    if border:
        shape.line.width = Pt(border_width)
        bc = border_color or "#D0D0D0"
        shape.line.color.rgb = hex_to_rgb(bc) if isinstance(bc, str) else bc
    else:
        shape.line.fill.background()
    return shape


def rect(slide, x, y, w, h, fill_color=None, border=False):
    """Add a plain rectangle. Useful for backgrounds, dividers.

    All position/size params are raw floats in inches (wrapped internally).

    Args:
        slide:      Slide object
        x, y:       Top-left position in inches (float, NOT Inches())
        w, h:       Width and height in inches (float, NOT Inches())
        fill_color: Hex string or RGBColor
        border:     True to show border line
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = (
            hex_to_rgb(fill_color) if isinstance(fill_color, str) else fill_color
        )
    if not border:
        shape.line.fill.background()
    return shape


def badge(
    slide,
    x,
    y,
    text,
    bg_color=None,
    text_color=None,
    size=10,
    padding_w=0.15,
    padding_h=0.05,
):
    """Add a small category badge/label. Returns (shape, textbox).

    x, y are raw floats in inches (wrapped internally).
    """
    font_h_in = size / 72
    w = len(text) * size * 0.65 / 72 + padding_w * 2
    h = font_h_in + padding_h * 2 + 0.05
    bg = bg_color or PALETTE["primary"]
    tc = text_color or PALETTE["white"]
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = hex_to_rgb(bg) if isinstance(bg, str) else bg
    pill.line.fill.background()
    tb = add_text(
        slide, x, y, w, h, text.upper(), size=size, bold=True, color=tc, align="center"
    )
    return pill, tb


def slide_background(slide, color):
    """Set solid background color for a slide.

    Args:
        slide: Slide object
        color: Hex string '#RRGGBB' or RGBColor
    """
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color) if isinstance(color, str) else color


def duplicate_slide(prs, slide_index):
    """Duplicate a slide by index. Returns the new slide."""
    template = prs.slides[slide_index]
    slide_layout = template.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in template.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)
    return new_slide


def inches_to_emu(inches):
    return int(inches * 914400)


def emu_to_inches(emu):
    return emu / 914400


def pt_to_inches(pt):
    return pt / 72


def text_height(num_lines, font_size_pt, line_spacing=1.05):
    """Calculate height in inches for N lines of text."""
    return (num_lines * font_size_pt * line_spacing) / 72
