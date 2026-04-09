import sys as _sys, os as _os

_SCRIPT_DIR = (
    _os.path.dirname(_os.path.abspath(__file__))
    if "__file__" in dir()
    else next(
        (p for p in _sys.path if _os.path.isfile(_os.path.join(p, "ppt_core.py"))),
        "/tmp/skills/ppt/scripts",
    )
)
if _SCRIPT_DIR not in _sys.path:
    _sys.path.insert(0, _SCRIPT_DIR)

"""
ppt_modifier.py — Surgical update operations on existing PPTX files.

Usage:
    from ppt_modifier import *
    replace_text('/tmp/deck.pptx', '/tmp/v2.pptx', 'DRAFT', 'FINAL')
    delete_slides('/tmp/v2.pptx', '/tmp/v3.pptx', [2, 5])
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy
from lxml import etree


def replace_text(
    filepath, output_path, find, replace, case_sensitive=True, slides=None
):
    prs = Presentation(filepath)
    count = 0
    for si, slide in enumerate(prs.slides, 1):
        if slides and si not in slides:
            continue
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if case_sensitive:
                            if find in run.text:
                                run.text = run.text.replace(find, replace)
                                count += 1
                        else:
                            if find.lower() in run.text.lower():
                                import re

                                run.text = re.sub(
                                    re.escape(find),
                                    replace,
                                    run.text,
                                    flags=re.IGNORECASE,
                                )
                                count += 1
    prs.save(output_path)
    return count


def replace_image(filepath, output_path, slide_num, shape_name, new_image_path):
    prs = Presentation(filepath)
    slide = prs.slides[slide_num - 1]
    for shape in slide.shapes:
        if shape.name == shape_name and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(new_image_path, left, top, width, height)
            prs.save(output_path)
            return True
    return False


def delete_slides(filepath, output_path, slide_nums):
    prs = Presentation(filepath)
    for si in sorted(slide_nums, reverse=True):
        if 1 <= si <= len(prs.slides):
            slide_elem = prs.slides._sldIdLst[si - 1]
            rId = slide_elem.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            prs.slides._sldIdLst.remove(slide_elem)
            try:
                prs.part.drop_rel(rId)
            except:
                pass
    prs.save(output_path)


def reorder_slides(filepath, output_path, new_order):
    prs = Presentation(filepath)
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    assert len(new_order) == len(items), (
        f"new_order length ({len(new_order)}) != slide count ({len(items)})"
    )
    assert sorted(new_order) == list(range(len(items))), (
        "new_order must contain each index exactly once"
    )
    reordered = [items[i] for i in new_order]
    for item in list(sldIdLst):
        sldIdLst.remove(item)
    for item in reordered:
        sldIdLst.append(item)
    prs.save(output_path)


def duplicate_slide(filepath, output_path, slide_num):
    prs = Presentation(filepath)
    source = prs.slides[slide_num - 1]
    new_slide = prs.slides.add_slide(source.slide_layout)
    for shape in source.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)
    try:
        src_bg = source.background._element
        if src_bg is not None:
            new_bg = copy.deepcopy(src_bg)
            new_slide.background._element.getparent().replace(
                new_slide.background._element, new_bg
            )
    except:
        pass
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    new_item = items[-1]
    sldIdLst.remove(new_item)
    sldIdLst.insert(slide_num, new_item)
    prs.save(output_path)
    return slide_num


def update_font_globally(filepath, output_path, old_font=None, new_font="Calibri"):
    prs = Presentation(filepath)
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if old_font is None or run.font.name == old_font:
                            run.font.name = new_font
                            count += 1
    prs.save(output_path)
    return count


def update_colors(filepath, output_path, color_map):
    prs = Presentation(filepath)
    count = 0
    cmap = {k.upper().lstrip("#"): v.upper().lstrip("#") for k, v in color_map.items()}
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.fill.type is not None:
                    c = str(shape.fill.fore_color.rgb).upper()
                    if c in cmap:
                        shape.fill.fore_color.rgb = RGBColor.from_string(cmap[c])
                        count += 1
            except:
                pass
            try:
                if shape.line.fill.type is not None:
                    c = str(shape.line.color.rgb).upper()
                    if c in cmap:
                        shape.line.color.rgb = RGBColor.from_string(cmap[c])
                        count += 1
            except:
                pass
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type is not None:
                                c = str(run.font.color.rgb).upper()
                                if c in cmap:
                                    run.font.color.rgb = RGBColor.from_string(cmap[c])
                                    count += 1
                        except:
                            pass
    prs.save(output_path)
    return count


def add_slide_numbers(
    filepath,
    output_path,
    start=1,
    x=12.2,
    y=7.0,
    size=10,
    color="#888888",
    font_name="Calibri",
):
    from ppt_core import add_text

    prs = Presentation(filepath)
    for si, slide in enumerate(prs.slides):
        add_text(
            slide,
            x,
            y,
            0.8,
            0.3,
            str(start + si),
            size=size,
            color=color,
            font_name=font_name,
            align="right",
        )
    prs.save(output_path)
