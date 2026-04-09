import sys as _sys, os as _os

_SCRIPT_DIR = (
    _os.path.dirname(_os.path.abspath(__file__))
    if "__file__" in dir()
    else next(
        (p for p in _sys.path if _os.path.isfile(_os.path.join(p, "ppt_core.py"))),
        "/tmp/skills/ppt/scripts",
    )
)
if _SCRIPT_DIR not in _sys.path:
    _sys.path.insert(0, _SCRIPT_DIR)

"""
ppt_images.py — Image sourcing, preparation, and safe placement.

IMPORTANT — COORDINATE CONVENTION:
    All position/size params (left, top, max_w, max_h) are raw floats in
    inches. They are wrapped with Inches() internally — do NOT pass Inches().

Usage:
    from ppt_images import *
    download_image('https://example.com/photo.jpg', '/tmp/images/photo.jpg')
    prepare_background('/tmp/images/photo.jpg', '/tmp/images/bg.jpg', darkness=0.3)
    add_full_bleed_bg(slide, '/tmp/images/bg.jpg')
    safe_add_picture(slide, '/tmp/images/chart.png', 1.0, 1.7, 10.0)
"""

from pptx.util import Inches
from PIL import Image, ImageEnhance, ImageFilter, ImageDraw
import requests
import os

SLIDE_W = 13.333
SLIDE_H = 7.5
TARGET_ASPECT = SLIDE_W / SLIDE_H
IMAGE_DIR = "/tmp/images"


def ensure_image_dir():
    os.makedirs(IMAGE_DIR, exist_ok=True)


def download_image(url, save_path=None, timeout=15):
    ensure_image_dir()
    if save_path is None:
        ext = os.path.splitext(url.split("?")[0])[-1] or ".jpg"
        save_path = os.path.join(IMAGE_DIR, f"download_{hash(url) % 100000}{ext}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }
        resp = requests.get(url, timeout=timeout, headers=headers, stream=True)
        resp.raise_for_status()
        with open(save_path, "wb") as f:
            for chunk in resp.iter_content(8192):
                f.write(chunk)
        img = Image.open(save_path)
        img.verify()
        return save_path
    except Exception as e:
        print(f"Download failed: {e}")
        if os.path.exists(save_path):
            os.remove(save_path)
        return None


def crop_to_16_9(input_path, output_path=None):
    if output_path is None:
        base, ext = os.path.splitext(input_path)
        output_path = f"{base}_16x9{ext}"
    img = Image.open(input_path).convert("RGB")
    w, h = img.size
    src_aspect = w / h
    if src_aspect > TARGET_ASPECT:
        new_w = int(h * TARGET_ASPECT)
        offset = (w - new_w) // 2
        img = img.crop((offset, 0, offset + new_w, h))
    else:
        new_h = int(w / TARGET_ASPECT)
        offset = (h - new_h) // 2
        img = img.crop((0, offset, w, offset + new_h))
    img.save(output_path, quality=95)
    return output_path


def prepare_background(
    input_path,
    output_path=None,
    darkness=0.3,
    tint_color=(15, 27, 45),
    tint_alpha=0.4,
    blur=2,
    target_size=(2667, 1500),
):
    if output_path is None:
        base, ext = os.path.splitext(input_path)
        output_path = f"{base}_bg{ext}"
    img = Image.open(input_path).convert("RGB")
    w, h = img.size
    src_aspect = w / h
    if src_aspect > TARGET_ASPECT:
        new_w = int(h * TARGET_ASPECT)
        offset = (w - new_w) // 2
        img = img.crop((offset, 0, offset + new_w, h))
    else:
        new_h = int(w / TARGET_ASPECT)
        offset = (h - new_h) // 2
        img = img.crop((0, offset, w, offset + new_h))
    img = img.resize(target_size, Image.LANCZOS)
    if blur > 0:
        img = img.filter(ImageFilter.GaussianBlur(radius=blur))
    img = ImageEnhance.Brightness(img).enhance(darkness)
    if tint_alpha > 0:
        tint = Image.new("RGB", img.size, tint_color)
        img = Image.blend(img, tint, alpha=tint_alpha)
    img.save(output_path, quality=95)
    return output_path


def create_gradient_background(
    output_path, color_top, color_bottom, width=2667, height=1500
):
    ensure_image_dir()
    img = Image.new("RGB", (width, height))
    draw = ImageDraw.Draw(img)
    for y in range(height):
        ratio = y / height
        r = int(color_top[0] + (color_bottom[0] - color_top[0]) * ratio)
        g = int(color_top[1] + (color_bottom[1] - color_top[1]) * ratio)
        b = int(color_top[2] + (color_bottom[2] - color_top[2]) * ratio)
        draw.line([(0, y), (width, y)], fill=(r, g, b))
    img.save(output_path, quality=95)
    return output_path


def add_full_bleed_bg(slide, image_path):
    """Place image as full-bleed background. Call FIRST so it's behind all shapes."""
    slide.shapes.add_picture(
        image_path, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H)
    )


def safe_add_picture(
    slide,
    img_path,
    left,
    top,
    max_w,
    max_h=None,
    slide_w=None,
    slide_h=None,
    footer_margin=0.5,
    # ── Aliases (accepted for convenience) ──
    width=None,
    height=None,
):
    """Place an image within a bounding box, preserving aspect ratio.

    All position/size params are raw floats in inches (wrapped internally).

    Args:
        slide:          Slide object
        img_path:       Path to the image file
        left, top:      Top-left position in inches (float, NOT Inches())
        max_w:          Maximum width in inches (float). Alias: width
        max_h:          Maximum height in inches (float, optional). Alias: height
        slide_w:        Override slide width (default: 13.333)
        slide_h:        Override slide height (default: 7.5)
        footer_margin:  Bottom margin to avoid (default: 0.5")

    Returns:
        (actual_width, actual_height) tuple in inches
    """
    # Resolve aliases
    if width is not None and max_w is None:
        max_w = width
    if height is not None and max_h is None:
        max_h = height

    sw = slide_w or SLIDE_W
    sh = slide_h or SLIDE_H
    if max_h is None:
        max_h = sh - top - footer_margin
    img = Image.open(img_path)
    aspect = img.size[0] / img.size[1]
    w = max_w
    h = w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    if left + w > sw - 0.3:
        w = sw - 0.3 - left
        h = w / aspect
    if top + h > sh - footer_margin:
        h = sh - footer_margin - top
        w = h * aspect
    slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(w), Inches(h))
    return (round(w, 2), round(h, 2))


def create_placeholder_image(
    width_px,
    height_px,
    output_path,
    text="Image Placeholder",
    bg_color=(220, 220, 230),
    text_color=(140, 140, 150),
):
    ensure_image_dir()
    img = Image.new("RGB", (width_px, height_px), bg_color)
    draw = ImageDraw.Draw(img)
    draw.line([(0, 0), (width_px, height_px)], fill=text_color, width=1)
    draw.line([(width_px, 0), (0, height_px)], fill=text_color, width=1)
    try:
        from PIL import ImageFont

        font = ImageFont.truetype(
            "/usr/share/fonts/google-droid-sans-fonts/DroidSans.ttf",
            max(16, min(width_px, height_px) // 15),
        )
    except:
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(
        ((width_px - tw) // 2, (height_px - th) // 2), text, fill=text_color, font=font
    )
    img.save(output_path, quality=90)
    return output_path


def composite_transparent_preview(chart_path, output_path, bg_color):
    """
    Composite a transparent-background chart PNG onto a solid color for
    accurate visual preview. Needed because LibreOffice (used by render_slides)
    fills transparent areas with white.

    The PPTX itself is correct — PowerPoint and Google Slides render
    transparency properly. This function is only for preview verification.

    Args:
        chart_path:  Path to the transparent chart PNG
        output_path: Where to save the composited preview
        bg_color:    Background color as (R, G, B) tuple or hex string.
                     Match this to your slide background for an accurate preview.
                     Use hex_to_rgb(c('slate', 950)) for dark-first decks.

    Returns:
        output_path

    Example:
        composite_transparent_preview(
            '/tmp/chart.png',
            '/tmp/chart_preview.png',
            hex_to_rgb(c('slate', 950))
        )
    """
    if isinstance(bg_color, str):
        bg_color = tuple(int(bg_color.lstrip("#")[i : i + 2], 16) for i in (0, 2, 4))

    chart = Image.open(chart_path).convert("RGBA")
    bg = Image.new("RGBA", chart.size, (*bg_color, 255))
    bg.paste(chart, (0, 0), chart)
    bg.convert("RGB").save(output_path, quality=95)
    return output_path
