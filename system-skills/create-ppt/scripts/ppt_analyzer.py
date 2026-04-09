import sys as _sys, os as _os

_SCRIPT_DIR = (
    _os.path.dirname(_os.path.abspath(__file__))
    if "__file__" in dir()
    else next(
        (p for p in _sys.path if _os.path.isfile(_os.path.join(p, "ppt_core.py"))),
        "/tmp/skills/ppt/scripts",
    )
)
if _SCRIPT_DIR not in _sys.path:
    _sys.path.insert(0, _SCRIPT_DIR)

"""
ppt_analyzer.py — Read, inspect, and analyze existing PPTX files.

Usage:
    from ppt_analyzer import *
    meta = read_metadata('/tmp/deck.pptx')
    report = analyze_template('/tmp/template.pptx')
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import Counter, defaultdict
import os, json


def _emu_to_in(emu):
    if emu is None:
        return 0.0
    return round(emu / 914400, 2)


def _color_str(color_obj):
    try:
        if color_obj and color_obj.rgb:
            return "#" + str(color_obj.rgb)
    except:
        pass
    return None


def read_metadata(filepath):
    prs = Presentation(filepath)
    cp = prs.core_properties
    return {
        "slide_count": len(prs.slides),
        "width_in": _emu_to_in(prs.slide_width),
        "height_in": _emu_to_in(prs.slide_height),
        "layout_count": len(prs.slide_layouts),
        "author": cp.author or "",
        "title": cp.title or "",
        "subject": cp.subject or "",
        "created": str(cp.created) if cp.created else "",
        "modified": str(cp.modified) if cp.modified else "",
        "last_modified_by": cp.last_modified_by or "",
    }


def read_all_text(filepath):
    prs = Presentation(filepath)
    results = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    results.append({"slide": si, "shape": shape.name, "text": text})
    return results


def read_slide_notes(filepath):
    prs = Presentation(filepath)
    notes = {}
    for si, slide in enumerate(prs.slides, 1):
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
            if text:
                notes[si] = text
    return notes


def extract_images(filepath, output_dir="/tmp/extracted_images"):
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(filepath)
    paths = []
    img_count = 0
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_count += 1
                    ext = os.path.splitext(shape.image.content_type.split("/")[-1])[-1]
                    if not ext.startswith("."):
                        ext = "." + ext
                    fname = f"slide{si:02d}_img{img_count:03d}{ext}"
                    out_path = os.path.join(output_dir, fname)
                    with open(out_path, "wb") as f:
                        f.write(shape.image.blob)
                    paths.append(out_path)
                except Exception as e:
                    paths.append(f"ERROR: slide {si}, shape {shape.name}: {e}")
    return paths


def inspect_shapes(filepath):
    prs = Presentation(filepath)
    shapes_data = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            info = {
                "slide": si,
                "name": shape.name,
                "type": str(shape.shape_type),
                "left_in": _emu_to_in(shape.left),
                "top_in": _emu_to_in(shape.top),
                "width_in": _emu_to_in(shape.width),
                "height_in": _emu_to_in(shape.height),
                "right_in": _emu_to_in((shape.left or 0) + (shape.width or 0)),
                "bottom_in": _emu_to_in((shape.top or 0) + (shape.height or 0)),
                "rotation": shape.rotation,
            }
            try:
                if shape.fill.type is not None:
                    info["fill_color"] = _color_str(shape.fill.fore_color)
                    info["fill_type"] = str(shape.fill.type)
            except:
                pass
            try:
                if shape.line.fill.type is not None:
                    info["line_color"] = _color_str(shape.line.color)
                    info["line_width_pt"] = (
                        round(shape.line.width.pt, 1) if shape.line.width else None
                    )
            except:
                pass
            if shape.has_text_frame:
                info["has_text"] = True
                info["text_preview"] = shape.text_frame.text[:80]
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            info["font_name"] = run.font.name
                            info["font_size_pt"] = (
                                run.font.size.pt if run.font.size else None
                            )
                            info["font_bold"] = run.font.bold
                            info["font_color"] = (
                                _color_str(run.font.color)
                                if run.font.color.type is not None
                                else None
                            )
                            break
                    if "font_name" in info:
                        break
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                info["is_picture"] = True
                try:
                    info["image_content_type"] = shape.image.content_type
                    info["image_size_bytes"] = len(shape.image.blob)
                except:
                    pass
            shapes_data.append(info)
    return shapes_data


def inspect_fonts(filepath):
    prs = Presentation(filepath)
    font_counter = Counter()
    font_details = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        fname = run.font.name or "Default"
                        fsize = run.font.size.pt if run.font.size else None
                        fbold = run.font.bold or False
                        fcolor = (
                            _color_str(run.font.color)
                            if run.font.color and run.font.color.type is not None
                            else None
                        )
                        key = (fname, fsize, fbold, fcolor)
                        font_counter[key] += 1
                        if key not in font_details:
                            font_details[key] = {
                                "font": fname,
                                "size_pt": fsize,
                                "bold": fbold,
                                "color": fcolor,
                            }
    results = []
    for key, count in font_counter.most_common():
        entry = font_details[key].copy()
        entry["count"] = count
        results.append(entry)
    return results


def extract_palette(filepath):
    prs = Presentation(filepath)
    fill_colors, text_colors, line_colors = set(), set(), set()
    for slide in prs.slides:
        try:
            bg = slide.background.fill
            if bg.type is not None:
                c = _color_str(bg.fore_color)
                if c:
                    fill_colors.add(c)
        except:
            pass
        for shape in slide.shapes:
            try:
                if shape.fill.type is not None:
                    c = _color_str(shape.fill.fore_color)
                    if c:
                        fill_colors.add(c)
            except:
                pass
            try:
                if shape.line.fill.type is not None:
                    c = _color_str(shape.line.color)
                    if c:
                        line_colors.add(c)
            except:
                pass
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type is not None:
                                c = _color_str(run.font.color)
                                if c:
                                    text_colors.add(c)
                        except:
                            pass
    all_colors = fill_colors | text_colors | line_colors
    return {
        "fill_colors": sorted(fill_colors),
        "text_colors": sorted(text_colors),
        "line_colors": sorted(line_colors),
        "all_colors": sorted(all_colors),
    }


def map_placeholders(filepath):
    prs = Presentation(filepath)
    layouts = []
    for layout in prs.slide_layouts:
        layout_info = {"name": layout.name, "placeholders": []}
        for ph in layout.placeholders:
            layout_info["placeholders"].append(
                {
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "name": ph.name,
                    "left_in": _emu_to_in(ph.left),
                    "top_in": _emu_to_in(ph.top),
                    "width_in": _emu_to_in(ph.width),
                    "height_in": _emu_to_in(ph.height),
                }
            )
        layouts.append(layout_info)
    return layouts


def analyze_template(filepath):
    report = {
        "metadata": read_metadata(filepath),
        "palette": extract_palette(filepath),
        "fonts": inspect_fonts(filepath),
        "layouts": map_placeholders(filepath),
        "text_content": read_all_text(filepath),
        "notes": read_slide_notes(filepath),
    }
    meta = report["metadata"]
    report["summary"] = (
        f"Presentation: {meta['slide_count']} slides, "
        f"{meta['width_in']}x{meta['height_in']}in, "
        f"{meta['layout_count']} layouts, "
        f"{len(report['palette']['all_colors'])} unique colors, "
        f"{len(report['fonts'])} font combinations"
    )
    return report
